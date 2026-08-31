"""Build a PPTX from the rendered deck, for import into Google Slides.

Each slide is placed as a full-bleed image so the typography and layout survive
the round-trip exactly - Google Slides re-flows native text boxes and would
break the design. The editable source stays deck.html.
"""
import subprocess
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches

PDF = Path("capitol-desk-deck.pdf")
OUT = Path("capitol-desk-deck.pptx")
TMP = Path("/tmp/pptx_slides")
TMP.mkdir(exist_ok=True)

for f in TMP.glob("*.png"):
    f.unlink()

# 200 dpi against a 960x540pt page gives a crisp 2667x1500 raster.
subprocess.run(
    ["magick", "-density", "200", str(PDF), "-background", "#0d0d0d",
     "-alpha", "remove", "-alpha", "off", "-quality", "95",
     str(TMP / "slide-%02d.png")],
    check=True,
)
pages = sorted(TMP.glob("slide-*.png"))
if not pages:
    raise SystemExit("no pages rendered")

prs = Presentation()
prs.slide_width = Inches(13.333)   # 16:9 widescreen
prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]       # completely empty layout

for png in pages:
    slide = prs.slides.add_slide(blank)
    slide.shapes.add_picture(
        str(png), 0, 0, width=prs.slide_width, height=prs.slide_height
    )

prs.save(OUT)
size = OUT.stat().st_size / 1_048_576
print(f"{OUT}: {len(pages)} slides, {size:.1f} MB")
